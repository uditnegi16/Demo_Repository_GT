import streamlit as st
import pandas as pd
import os
from module.data_ingestion import DataIngestor
from module.data_processing import DataProcessor
from module.ai_insight import GeminiInsights
from module.visualization import DataVisualizer
from module.report_pdf import PDFReportGenerator
from module.report_pptx import PowerPointReportGenerator

# Page configuration
st.set_page_config(
    page_title="AdTech Report Generator",
    page_icon="📊",
    layout="wide"
)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1E3A8A;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        color: #374151;
        text-align: center;
        margin-bottom: 2rem;
    }
    .stButton>button {
        width: 100%;
        margin-top: 10px;
    }
    .success-box {
        padding: 1rem;
        background-color: #D1FAE5;
        border-radius: 0.5rem;
        border-left: 4px solid #10B981;
    }
</style>
""", unsafe_allow_html=True)

# App title
st.markdown('<h1 class="main-header">🚀 AdTech Report Automation</h1>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">Upload data → AI Insights → Download Reports (PDF/PPT)</p>', unsafe_allow_html=True)

# Initialize session state
if 'data' not in st.session_state:
    st.session_state.data = None
if 'ingestor' not in st.session_state:
    st.session_state.ingestor = DataIngestor()
if 'processor' not in st.session_state:
    st.session_state.processor = None
if 'ai_insights' not in st.session_state:
    st.session_state.ai_insights = GeminiInsights()
if 'insights_generated' not in st.session_state:
    st.session_state.insights_generated = False
if 'cleaned_data' not in st.session_state:
    st.session_state.cleaned_data = None

# Sidebar - Data Upload
with st.sidebar:
    st.header("📁 Data Source")
    
    data_source = st.radio(
        "Choose your data source:",
        ["CSV File", "Excel File", "SQL Database", "Use Sample Data"],
        index=0
    )
    
    if data_source in ["CSV File", "Excel File"]:
        file_type = "csv" if data_source == "CSV File" else "xlsx"
        uploaded_file = st.file_uploader(
            f"Upload {data_source}",
            type=[file_type, 'xls'],
            help=f"Upload your {data_source.lower()}"
        )
        
        if uploaded_file is not None:
            if data_source == "CSV File":
                st.session_state.data = st.session_state.ingestor.ingest_csv(uploaded_file)
            else:
                st.session_state.data = st.session_state.ingestor.ingest_excel(uploaded_file)
            
            if st.session_state.data is not None:
                st.session_state.processor = DataProcessor(st.session_state.data)
                st.session_state.cleaned_data = None
    
    elif data_source == "SQL Database":
        with st.form("sql_connection"):
            st.subheader("Database Connection")
            db_type = st.selectbox("Database Type", ["MySQL", "PostgreSQL", "SQLite"])
            
            col1, col2 = st.columns(2)
            with col1:
                host = st.text_input("Host", "localhost")
                database = st.text_input("Database", "mydb")
            with col2:
                port = st.text_input("Port", "3306" if db_type == "MySQL" else "5432")
            
            username = st.text_input("Username", "root")
            password = st.text_input("Password", type="password")
            query = st.text_area("SQL Query", "SELECT * FROM ad_clicks LIMIT 1000")
            
            if st.form_submit_button("Connect to Database"):
                with st.spinner("Connecting to database..."):
                    if db_type == "MySQL":
                        conn_string = f"mysql+pymysql://{username}:{password}@{host}:{port}/{database}"
                    elif db_type == "PostgreSQL":
                        conn_string = f"postgresql://{username}:{password}@{host}:{port}/{database}"
                    else:
                        conn_string = f"sqlite:///{database}"
                    
                    st.session_state.data = st.session_state.ingestor.ingest_sql(conn_string, query)
                    if st.session_state.data is not None:
                        st.session_state.processor = DataProcessor(st.session_state.data)
    
    else:  # Sample Data
        if st.button("📊 Load Sample Dataset", use_container_width=True):
            # Placeholder for sample data
            st.info("Sample data feature will be implemented in Phase 2")
            # We'll load the Kaggle dataset here

# Main Content Area
if st.session_state.data is not None:
    # Data Preview Section
    st.header("📊 Data Preview")
    
    tab1, tab2, tab3, tab4 = st.tabs(["Raw Data", "Data Info", "Basic Stats", "Data Cleaning"])
    
    with tab1:
        st.subheader("First 10 Rows")
        st.dataframe(st.session_state.data.head(10), use_container_width=True)
        
        # Show total rows
        total_rows = len(st.session_state.data)
        st.caption(f"Total rows: {total_rows:,}")
    
    with tab2:
        info = st.session_state.ingestor.get_data_info()
        if info:
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Rows", f"{info['shape'][0]:,}")
                st.metric("Columns", info['shape'][1])
                st.metric("Memory", f"{info['memory_usage']:.2f} MB")
            
            with col2:
                st.write("**Columns List:**")
                for col in info['columns']:
                    st.code(col)
    
    with tab3:
        st.subheader("Statistical Summary")
        st.dataframe(st.session_state.data.describe(), use_container_width=True)
    
    with tab4:
        st.subheader("Data Cleaning")
        
        if st.button("🧹 Clean Data", type="secondary"):
            with st.spinner("Cleaning data..."):
                st.session_state.cleaned_data = st.session_state.processor.clean_data()
        
        if st.session_state.cleaned_data is not None:
            st.success("✅ Data cleaned successfully!")
            
            # Show cleaning results
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Original Rows", len(st.session_state.data))
                st.metric("Cleaned Rows", len(st.session_state.cleaned_data))
            with col2:
                date_cols = st.session_state.processor.detect_date_columns()
                st.write(f"📅 Date columns detected: {len(date_cols)}")
                if date_cols:
                    st.write(date_cols)
    
    # Divider
    st.markdown("---")
    
    # AI Analysis Section - UPDATED
    st.markdown("---")
    st.header("🤖 AI-Powered Analysis")
    
    col1, col2 = st.columns([3, 1])
    with col1:
        if st.button("🔍 Generate Comprehensive AI Insights", type="primary", use_container_width=True):
            with st.spinner("🤔 Analyzing data with Gemini AI..."):
                # Get AI insights
                ai = GeminiInsights()
                ai_summary = ai.analyze_adtech_data(st.session_state.data)
                
                # Generate visualizations
                visualizer = DataVisualizer(st.session_state.data)
                charts = visualizer.create_summary_charts()
                adtech_charts = visualizer.create_adtech_specific_charts()
                
                # Store in session state
                st.session_state.ai_summary = ai_summary
                st.session_state.charts = charts
                st.session_state.adtech_charts = adtech_charts
                st.session_state.insights_generated = True
                st.session_state.visualizations_ready = True
    
    with col2:
        if st.button("🔄 Clear Analysis", type="secondary"):
            st.session_state.insights_generated = False
            st.session_state.visualizations_ready = False
            st.rerun()
    
    # Display AI insights if generated
    if st.session_state.insights_generated and 'ai_summary' in st.session_state:
        st.subheader("AI-Generated Insights")
        
        with st.expander("📋 View Analysis", expanded=True):
            st.markdown(st.session_state.ai_summary)
        
        # Report Generation Options
        st.markdown("---")
        st.header("📄 Report Generation")

        col1, col2 = st.columns(2)

        with col1:
            if st.button("📥 Download PDF Report", use_container_width=True, type="primary"):
                with st.spinner("Generating PDF report..."):
                    try:
                        from module.report_pdf import PDFReportGenerator
                        
                        # Create PDF
                        pdf_gen = PDFReportGenerator()
                        
                        # Get data and insights
                        data = st.session_state.data
                        insights = st.session_state.ai_summary if 'ai_summary' in st.session_state else "AI insights not generated"
                        charts = st.session_state.charts if 'charts' in st.session_state else {}
                        
                        # Generate PDF
                        pdf_filename = pdf_gen.generate_simple_report(data, insights)
                        
                        # Read the file
                        with open(pdf_filename, "rb") as pdf_file:
                            pdf_bytes = pdf_file.read()
                        
                        # Create download button
                        st.download_button(
                            label="⬇️ Click to Download PDF",
                            data=pdf_bytes,
                            file_name=f"adtech_report_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pdf",
                            mime="application/pdf",
                            key="pdf_download",
                            use_container_width=True
                        )
                        st.success("✅ PDF report generated!")
                        
                    except Exception as e:
                        st.error(f"PDF Error: {str(e)}")
                        # Fallback: Create simple PDF
                        import io
                        from reportlab.pdfgen import canvas
                        from reportlab.lib.pagesizes import letter
                        
                        buffer = io.BytesIO()
                        c = canvas.Canvas(buffer, pagesize=letter)
                        c.drawString(100, 750, "AdTech Performance Report")
                        c.drawString(100, 730, f"Generated: {pd.Timestamp.now()}")
                        c.drawString(100, 710, f"Rows: {len(st.session_state.data)}")
                        c.drawString(100, 690, "Charts and insights included in full version")
                        c.save()
                        
                        st.download_button(
                            label="⬇️ Download Basic PDF",
                            data=buffer.getvalue(),
                            file_name="adtech_report_basic.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )

        with col2:
            if st.button("📊 Download PowerPoint", use_container_width=True, type="primary"):
                with st.spinner("Generating PowerPoint presentation..."):
                    try:
                        from module.report_pptx import PowerPointReportGenerator
                        
                        # Create PowerPoint
                        pptx_gen = PowerPointReportGenerator()
                        
                        # Get data and insights
                        data = st.session_state.data
                        insights = st.session_state.ai_summary if 'ai_summary' in st.session_state else "AI insights not generated"
                        
                        # Generate PPTX
                        pptx_filename = pptx_gen.generate_simple_presentation(data, insights)
                        
                        # Read the file
                        with open(pptx_filename, "rb") as pptx_file:
                            pptx_bytes = pptx_file.read()
                        
                        # Create download button
                        st.download_button(
                            label="⬇️ Click to Download PowerPoint",
                            data=pptx_bytes,
                            file_name=f"adtech_presentation_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="pptx_download",
                            use_container_width=True
                        )
                        st.success("✅ PowerPoint presentation generated!")
                        
                    except Exception as e:
                        st.error(f"PowerPoint Error: {str(e)}")
                        # Create a simple PPTX as fallback
                        from pptx import Presentation
                        import io
                        
                        prs = Presentation()
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = "AdTech Report"
                        slide.placeholders[1].text = "Generated by TrendSpotter"
                        
                        buffer = io.BytesIO()
                        prs.save(buffer)
                        
                        st.download_button(
                            label="⬇️ Download Basic PowerPoint",
                            data=buffer.getvalue(),
                            file_name="adtech_presentation_basic.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
    
else:
    # Welcome screen when no data is loaded
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        st.image("https://cdn-icons-png.flaticon.com/512/3067/3067256.png", width=150)
        st.markdown("""
        ### Welcome to AdTech Report Automation
        
        **Get started in 3 simple steps:**
        1. **Upload** your data (CSV, Excel, or SQL)
        2. **Analyze** with AI-powered insights
        3. **Download** professional reports (PDF/PPT)
        
        👈 **Select your data source from the sidebar**
        """)
        
        # Quick tips
        with st.expander("💡 Tips for best results"):
            st.markdown("""
            - **CSV Files**: Ensure consistent column names
            - **Excel Files**: First sheet will be used by default
            - **SQL Databases**: Use SELECT queries with LIMIT for large datasets
            - **Data should include**: Timestamps, metrics (clicks, conversions), and categorical data
            """)

# Footer
st.markdown("---")
st.markdown("### 🚀 Hackathon Progress")
progress_cols = st.columns(4)

with progress_cols[0]:
    st.metric("Phase", "3/3", "✅")
with progress_cols[1]:
    st.metric("Status", "Report Generation", "Complete")
with progress_cols[2]:
    st.metric("Features", "PDF & PPT", "Ready")
with progress_cols[3]:
    st.metric("Time", "12:45", "On Track")

st.caption("Made for GroundTruth AI Fellowship Hackathon | TrendSpotter v1.0 Complete")
