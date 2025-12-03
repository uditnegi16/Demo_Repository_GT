from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from datetime import datetime
import io

class PowerPointReportGenerator:
    def __init__(self):
        self.prs = None
        
    def _truncate_text_for_pptx(self, text, max_lines=10, max_chars_per_line=80):
    # """Truncate text to fit PowerPoint slides without overflow"""
        lines = text.split('\n')
        truncated_lines = []
        
        for line in lines[:max_lines]:  # Limit to max_lines
            if len(line) > max_chars_per_line:
                # Split long lines
                words = line.split()
                current_line = ""
                for word in words:
                    if len(current_line) + len(word) + 1 <= max_chars_per_line:
                        current_line += f" {word}" if current_line else word
                    else:
                        truncated_lines.append(current_line)
                        current_line = word
                if current_line:
                    truncated_lines.append(current_line)
            else:
                truncated_lines.append(line)
        
        # Add ellipsis if text was truncated
        if len(lines) > max_lines or any(len(line) > max_chars_per_line for line in lines):
            if truncated_lines and not truncated_lines[-1].endswith('...'):
                truncated_lines[-1] = truncated_lines[-1] + "..."
        
        return '\n'.join(truncated_lines)

    def _add_text_to_shape(self, shape, text):
        # """Safely add text to PowerPoint shape with proper formatting"""
        # Clear existing text
        shape.text = ""
        
        # Truncate if needed
        truncated_text = self._truncate_text_for_pptx(text)
        
        # Add text with proper paragraph formatting
        text_frame = shape.text_frame
        text_frame.clear()  # Clear any existing paragraphs
        
        # Split into paragraphs
        paragraphs = truncated_text.split('\n')
        
        for i, para_text in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = para_text
            # Adjust font size based on text length
            if len(para_text) > 60:
                p.font.size = Pt(10)
            else:
                p.font.size = Pt(12)
                
    def generate_simple_presentation(self, df, ai_insights, filename="adtech_presentation.pptx"):
        """Generate a simple PowerPoint presentation"""
        
        # Create presentation
        prs = Presentation()
        
        # Slide 1: Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "AdTech Performance Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}\nTrendSpotter Automated Insights"
        
        # Slide 2: Executive Summary
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        
        # Add AI insights summary (first 500 chars)
        summary = ai_insights[:500] + "..." if len(ai_insights) > 500 else ai_insights
        content.text = self._truncate_text_for_pptx(f"Key Insights:\n\n{summary}")
        
        # Slide 3: Dataset Overview
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Dataset Overview"
        
        dataset_info = f"""
        • Total Rows: {len(df):,}
        • Total Columns: {len(df.columns)}
        • Data Types:
        """
        
        # Add column info
        numeric_cols = df.select_dtypes(include=['number']).columns
        cat_cols = df.select_dtypes(include=['object']).columns
        
        dataset_info += f"\n• Numeric Columns ({len(numeric_cols)}):"
        for col in numeric_cols[:5]:  # First 5 only
            dataset_info += f"\n   - {col}"
        
        if len(numeric_cols) > 5:
            dataset_info += f"\n   ... and {len(numeric_cols) - 5} more"
        
        dataset_info += f"\n\n• Categorical Columns ({len(cat_cols)}):"
        for col in cat_cols[:3]:  # First 3 only
            dataset_info += f"\n   - {col}"
        
        if len(cat_cols) > 3:
            dataset_info += f"\n   ... and {len(cat_cols) - 3} more"
        
        content.text = dataset_info
        
        # Slide 4: Key Metrics
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Metrics"
        
        metrics_text = "Summary Statistics:\n\n"
        if len(numeric_cols) > 0:
            for col in numeric_cols[:6]:  # First 6 numeric columns
                metrics_text += f"• {col}:\n"
                metrics_text += f"  Mean: {df[col].mean():.2f}\n"
                metrics_text += f"  Min: {df[col].min():.2f} | Max: {df[col].max():.2f}\n"
                metrics_text += f"  Std Dev: {df[col].std():.2f}\n\n"
        
        content.text = self._truncate_text_for_pptx(metrics_text, max_lines=15)
        
        # Slide 5: Recommendations
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Actionable Recommendations"
        
        # Extract recommendations from AI insights
        rec_text = "Based on AI Analysis:\n\n"
        
        # Look for recommendations in AI text
        lines = ai_insights.split('\n')
        rec_lines = [line for line in lines if 'recommend' in line.lower() or 'suggest' in line.lower()]
        
        if rec_lines:
            for i, line in enumerate(rec_lines[:5], 1):  # First 5 recommendations
                rec_text += f"{i}. {line.strip()}\n\n"
        else:
            rec_text += "1. Optimize campaign targeting based on performance metrics\n"
            rec_text += "2. Consider A/B testing for underperforming ad creatives\n"
            rec_text += "3. Reallocate budget to high-conversion channels\n"
            rec_text += "4. Implement real-time monitoring for anomaly detection\n"
            rec_text += "5. Schedule regular performance reviews\n"
        
        content.text = self._truncate_text_for_pptx(rec_text, max_lines=8)
        
        # Slide 6: Thank You
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add text box
        from pptx.util import Inches
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "Thank You"
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = "Generated by TrendSpotter\nAutomated Insights Engine"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        
        # Save presentation
        prs.save(filename)
        
        return filename